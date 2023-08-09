import os
from fastapi import FastAPI, File, UploadFile, Form 
from typing import Annotated
from pydantic import BaseModel
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile

class Prompt(BaseModel):
    prompt: str

app = FastAPI()

@app.get("/health")
async def health():
    return {"Message": "Up and Running"}

@app.post("/")
async def create_item(res: Prompt):
    return res


@app.post("/upload/")
async def upload(prompt: Annotated[str, Form()], file: Annotated[UploadFile, File()]):
    temp_dir = tempfile.mkdtemp()
    file_path = os.path.join(temp_dir, file.filename)

    with open(file_path, "wb") as f:
        f.write(file.file.read())

    if (file.filename.endswith(".pdf")):
        print("It's a pdf file")
        reader = PyPDF2.PdfReader(file.file)
        for page in reader.pages:
            content = page.extract_text()
            print(content)

    elif (file.filename.endswith(".docx") or file.filename.endswith(".doc")):
        print("It's doc file")
        doc = Document(file_path)
        content = ""
        for para in doc.paragraphs:
            content += para.text
        print(content)
            
    elif (file.filename.endswith(".ppt") or file.filename.endswith(".pptx")):
        print("It's a ppt file")
        content = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    content +=  shape.text + "\n"
        print(content)  

    elif (file.filename.endswith(".txt")):
        print("It's a txt file")
        content = ""
        with open(file_path, 'r', encoding='utf-8') as f:
           content = f.read() 
        print(content)

    return {"prompt": prompt, "file": file.filename}